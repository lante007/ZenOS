#!/usr/bin/env python3
"""Extract slide text and speaker notes from a .pptx file, printed as
plain text with slide separators. Used as a preprocessing step by
src/text-extractor.js — output is fed straight into the same
length-based quality gate used for PDFs."""
import sys
from pptx import Presentation


def main(path):
    prs = Presentation(path)
    blocks = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in para.runs)
                    if text.strip():
                        texts.append(text.strip())
        slide_text = ' '.join(texts)

        notes = ''
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        block = '--- SLIDE {} ---\n{}'.format(i, slide_text)
        if notes:
            block += '\n[NOTES: {}]'.format(notes)
        blocks.append(block)

    print('\n\n'.join(blocks))


if __name__ == '__main__':
    main(sys.argv[1])
